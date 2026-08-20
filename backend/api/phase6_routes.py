"""Phase 6 API routes: File uploads, code execution, and multi-modal analysis.

Security features:
- Magic byte detection for file type verification
- Virus scanning (optional, requires ClamAV)
- Rate limiting per IP
- File metadata storage
- No file size or type restrictions (personal agent)
"""

from fastapi import APIRouter, UploadFile, File, HTTPException, Depends, Request
from fastapi.responses import FileResponse
from fastapi.security import APIKeyHeader
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
from pathlib import Path
import shutil
import uuid
import asyncio
import json
from datetime import datetime, timedelta
from collections import defaultdict
import hashlib

from app.utils.logger import app_logger
from app.tools.universal_filesystem import UniversalFilesystem
from app.tools.disposable_sandbox import DisposableSandbox
from app.tools.vision_analyzer import VisionAnalyzerTool
from app.tools.ocr_reader import OCRReaderTool

router = APIRouter(prefix="/api", tags=["phase6"])

# Configuration
UPLOAD_DIR = Path("./uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

METADATA_FILE = UPLOAD_DIR / ".metadata.json"

# Rate limiting configuration
RATE_LIMIT_REQUESTS = 100  # requests per window
RATE_LIMIT_WINDOW = timedelta(minutes=1)
rate_limit_store: Dict[str, List[datetime]] = defaultdict(list)

# Optional: ClamAV configuration
CLAMAV_ENABLED = False  # Set to True if ClamAV is installed
CLAMAV_SOCKET = "/var/run/clamav/clamd.ctl"


# ============================================================================
# Security Utilities
# ============================================================================

def detect_file_type(content: bytes, filename: Optional[str] = None) -> Dict[str, Any]:
    """
    Detect file type using magic bytes and filename extension.
    
    Returns:
        Dict with 'mime_type', 'category', 'extension', 'confidence'
    """
    # Magic byte signatures
    signatures = {
        # Images
        b'\xFF\xD8\xFF': ('image/jpeg', 'image', '.jpg'),
        b'\x89PNG\r\n\x1a\n': ('image/png', 'image', '.png'),
        b'GIF87a': ('image/gif', 'image', '.gif'),
        b'GIF89a': ('image/gif', 'image', '.gif'),
        b'RIFF': ('image/webp', 'image', '.webp'),  # Also used by WAV/AVI
        
        # Documents
        b'%PDF': ('application/pdf', 'document', '.pdf'),
        b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1': ('application/msword', 'document', '.doc'),
        b'PK\x03\x04': ('application/zip', 'archive', '.zip'),  # Also DOCX, XLSX, etc.
        
        # Video
        b'\x00\x00\x00\x18ftypmp4': ('video/mp4', 'video', '.mp4'),
        b'\x00\x00\x00\x1Cftypmp4': ('video/mp4', 'video', '.mp4'),
        b'\x00\x00\x00\x20ftypisom': ('video/mp4', 'video', '.mp4'),
        b'\x1A\x45\xDF\xA3': ('video/webm', 'video', '.webm'),
        b'moov': ('video/quicktime', 'video', '.mov'),
        
        # Audio
        b'ID3': ('audio/mpeg', 'audio', '.mp3'),
        b'\xFF\xFB': ('audio/mpeg', 'audio', '.mp3'),
        b'\xFF\xF3': ('audio/mpeg', 'audio', '.mp3'),
        b'\xFF\xF2': ('audio/mpeg', 'audio', '.mp3'),
        b'RIFF': ('audio/wav', 'audio', '.wav'),
        b'OggS': ('audio/ogg', 'audio', '.ogg'),
        b'fLaC': ('audio/flac', 'audio', '.flac'),
        
        # Archives
        b'PK\x03\x04': ('application/zip', 'archive', '.zip'),
        b'Rar!\x1A\x07': ('application/x-rar-compressed', 'archive', '.rar'),
        b'\x1F\x8B': ('application/gzip', 'archive', '.gz'),
        b'7z\xBC\xAF\x27\x1C': ('application/x-7z-compressed', 'archive', '.7z'),
        
        # Code/Text - check for Python shebang first
        b'#!/usr/bin/env python3': ('text/x-python', 'code', '.py'),
        b'#!/usr/bin/env python': ('text/x-python', 'code', '.py'),
        b'#!/usr/bin/python3': ('text/x-python', 'code', '.py'),
        b'#!/usr/bin/python': ('text/x-python', 'code', '.py'),
        b'#!/': ('text/x-shellscript', 'code', '.sh'),
        b'<?php': ('text/x-php', 'code', '.php'),
    }
    
    # Check magic bytes (first 32 bytes to accommodate longer shebangs)
    header = content[:32]
    detected_type = None
    detected_category = None
    detected_ext = None
    confidence = 'low'
    
    for signature, (mime_type, category, ext) in signatures.items():
        if header.startswith(signature):
            detected_type = mime_type
            detected_category = category
            detected_ext = ext
            confidence = 'high'
            break
    
    # Fallback to filename extension
    if not detected_type and filename:
        ext = Path(filename).suffix.lower()
        ext_to_mime = {
            '.jpg': ('image/jpeg', 'image'),
            '.jpeg': ('image/jpeg', 'image'),
            '.png': ('image/png', 'image'),
            '.gif': ('image/gif', 'image'),
            '.webp': ('image/webp', 'image'),
            '.pdf': ('application/pdf', 'document'),
            '.doc': ('application/msword', 'document'),
            '.docx': ('application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'document'),
            '.xls': ('application/vnd.ms-excel', 'document'),
            '.xlsx': ('application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'document'),
            '.ppt': ('application/vnd.ms-powerpoint', 'document'),
            '.pptx': ('application/vnd.openxmlformats-officedocument.presentationml.presentation', 'document'),
            '.txt': ('text/plain', 'text'),
            '.md': ('text/markdown', 'text'),
            '.py': ('text/x-python', 'code'),
            '.js': ('text/javascript', 'code'),
            '.ts': ('text/typescript', 'code'),
            '.html': ('text/html', 'code'),
            '.css': ('text/css', 'code'),
            '.json': ('application/json', 'document'),
            '.xml': ('application/xml', 'document'),
            '.yaml': ('text/yaml', 'document'),
            '.yml': ('text/yaml', 'document'),
            '.mp4': ('video/mp4', 'video'),
            '.webm': ('video/webm', 'video'),
            '.mov': ('video/quicktime', 'video'),
            '.avi': ('video/x-msvideo', 'video'),
            '.mp3': ('audio/mpeg', 'audio'),
            '.wav': ('audio/wav', 'audio'),
            '.ogg': ('audio/ogg', 'audio'),
            '.flac': ('audio/flac', 'audio'),
            '.zip': ('application/zip', 'archive'),
            '.tar': ('application/x-tar', 'archive'),
            '.gz': ('application/gzip', 'archive'),
            '.rar': ('application/x-rar-compressed', 'archive'),
            '.7z': ('application/x-7z-compressed', 'archive'),
        }
        
        if ext in ext_to_mime:
            detected_type, detected_category = ext_to_mime[ext]
            detected_ext = ext
            confidence = 'medium'
    
    # Default fallback
    if not detected_type:
        detected_type = 'application/octet-stream'
        detected_category = 'binary'
        detected_ext = Path(filename).suffix.lower() if filename else ''
        confidence = 'low'
    
    return {
        'mime_type': detected_type,
        'category': detected_category,
        'extension': detected_ext,
        'confidence': confidence,
    }


def calculate_file_hash(content: bytes) -> str:
    """Calculate SHA-256 hash of file content."""
    return hashlib.sha256(content).hexdigest()


async def scan_for_viruses(file_path: Path) -> Dict[str, Any]:
    """
    Scan file for viruses using ClamAV (if available).
    
    Returns:
        Dict with 'clean' (bool), 'virus_name' (optional), 'error' (optional)
    """
    if not CLAMAV_ENABLED:
        return {'clean': True, 'skipped': True}
    
    try:
        import pyclamd
        
        cd = pyclamd.ClamdUnixSocket(CLAMAV_SOCKET)
        
        # Test connection
        if not cd.ping():
            return {'clean': True, 'skipped': True, 'error': 'ClamAV not responding'}
        
        # Scan file
        result = cd.scan_file(str(file_path))
        
        if result is None:
            return {'clean': True}
        else:
            virus_name = list(result.values())[0][1]
            return {'clean': False, 'virus_name': virus_name}
    
    except Exception as e:
        app_logger.warning(f"Virus scan failed: {e}")
        return {'clean': True, 'skipped': True, 'error': str(e)}


def check_rate_limit(ip: str) -> bool:
    """
    Check if IP has exceeded rate limit.
    
    Returns:
        True if request is allowed, False if rate limited
    """
    now = datetime.utcnow()
    window_start = now - RATE_LIMIT_WINDOW
    
    # Clean old entries
    rate_limit_store[ip] = [
        timestamp for timestamp in rate_limit_store[ip]
        if timestamp > window_start
    ]
    
    # Check limit
    if len(rate_limit_store[ip]) >= RATE_LIMIT_REQUESTS:
        return False
    
    # Record this request
    rate_limit_store[ip].append(now)
    return True


def load_metadata() -> Dict[str, Any]:
    """Load file metadata from disk."""
    if METADATA_FILE.exists():
        try:
            return json.loads(METADATA_FILE.read_text())
        except Exception as e:
            app_logger.warning(f"Failed to load metadata: {e}")
    return {}


def save_metadata(metadata: Dict[str, Any]):
    """Save file metadata to disk."""
    try:
        METADATA_FILE.write_text(json.dumps(metadata, indent=2))
    except Exception as e:
        app_logger.error(f"Failed to save metadata: {e}")


# ============================================================================
# File Upload & Management (Phase 6a)
# ============================================================================

class FileUploadResponse(BaseModel):
    id: str
    name: str
    path: str
    size: int
    type: str
    category: str
    hash: str
    uploadedAt: str
    conversationId: Optional[str] = None


class FileListResponse(BaseModel):
    files: List[FileUploadResponse]
    total: int


@router.post("/files/upload", response_model=FileUploadResponse)
async def upload_file(
    request: Request,
    file: UploadFile = File(...),
    conversationId: Optional[str] = None
):
    """
    Upload a file to the server.
    
    No size or type restrictions. All files are accepted and analyzed.
    """
    # Rate limiting
    client_ip = request.client.host if request.client else "unknown"
    if not check_rate_limit(client_ip):
        raise HTTPException(
            status_code=429,
            detail=f"Rate limit exceeded. Max {RATE_LIMIT_REQUESTS} requests per minute."
        )
    
    # Read file content
    content = await file.read()
    file_size = len(content)
    
    app_logger.info(f"Uploading file: {file.filename} ({file_size} bytes)")
    
    # Detect file type using magic bytes
    file_type_info = detect_file_type(content, file.filename)
    
    # Calculate hash
    file_hash = calculate_file_hash(content)
    
    # Generate unique file ID and path
    file_id = str(uuid.uuid4())
    file_ext = file_type_info['extension'] or (Path(file.filename).suffix if file.filename else "")
    file_name = f"{file_id}{file_ext}"
    file_path = UPLOAD_DIR / file_name
    
    # Save file
    try:
        with open(file_path, "wb") as f:
            f.write(content)
        app_logger.info(f"File saved: {file_name}")
    except Exception as e:
        app_logger.error(f"Failed to save file: {e}")
        raise HTTPException(status_code=500, detail="Failed to save file")
    
    # Virus scan (optional)
    scan_result = await scan_for_viruses(file_path)
    if not scan_result.get('clean'):
        # Delete infected file
        file_path.unlink()
        virus_name = scan_result.get('virus_name', 'unknown')
        app_logger.warning(f"Virus detected in {file.filename}: {virus_name}")
        raise HTTPException(
            status_code=400,
            detail=f"File rejected: virus detected ({virus_name})"
        )
    
    # Store metadata
    metadata = load_metadata()
    metadata[file_id] = {
        'id': file_id,
        'name': file.filename or file_name,
        'path': str(file_path),
        'size': file_size,
        'type': file_type_info['mime_type'],
        'category': file_type_info['category'],
        'extension': file_ext,
        'hash': file_hash,
        'uploadedAt': datetime.utcnow().isoformat(),
        'conversationId': conversationId,
        'type_confidence': file_type_info['confidence'],
        'virus_scan': scan_result,
    }
    save_metadata(metadata)
    
    return FileUploadResponse(
        id=file_id,
        name=file.filename or file_name,
        path=str(file_path),
        size=file_size,
        type=file_type_info['mime_type'],
        category=file_type_info['category'],
        hash=file_hash,
        uploadedAt=metadata[file_id]['uploadedAt'],
        conversationId=conversationId
    )


@router.get("/files/{file_id}")
async def download_file(file_id: str):
    """Download a file by ID."""
    metadata = load_metadata()
    
    if file_id not in metadata:
        raise HTTPException(status_code=404, detail="File not found")
    
    file_info = metadata[file_id]
    file_path = Path(file_info['path'])
    
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found on disk")
    
    return FileResponse(
        path=str(file_path),
        filename=file_info['name'],
        media_type=file_info['type']
    )


@router.get("/files", response_model=FileListResponse)
async def list_files(conversationId: Optional[str] = None):
    """List uploaded files, optionally filtered by conversation."""
    metadata = load_metadata()
    
    files = []
    for file_id, file_info in metadata.items():
        if conversationId and file_info.get('conversationId') != conversationId:
            continue
        
        files.append(FileUploadResponse(
            id=file_info['id'],
            name=file_info['name'],
            path=file_info['path'],
            size=file_info['size'],
            type=file_info['type'],
            category=file_info['category'],
            hash=file_info['hash'],
            uploadedAt=file_info['uploadedAt'],
            conversationId=file_info.get('conversationId')
        ))
    
    return FileListResponse(files=files, total=len(files))


@router.delete("/files/{file_id}")
async def delete_file(file_id: str):
    """Delete a file by ID."""
    metadata = load_metadata()
    
    if file_id not in metadata:
        raise HTTPException(status_code=404, detail="File not found")
    
    file_info = metadata[file_id]
    file_path = Path(file_info['path'])
    
    # Delete file from disk
    if file_path.exists():
        file_path.unlink()
        app_logger.info(f"File deleted: {file_info['name']}")
    
    # Remove from metadata
    del metadata[file_id]
    save_metadata(metadata)
    
    return {"success": True, "message": "File deleted"}


# ============================================================================
# Code Execution (Phase 6b)
# ============================================================================

class CodeExecutionRequest(BaseModel):
    code: str
    language: str
    timeout: int = 30  # seconds (hard-capped at MAX_CODE_TIMEOUT_SECONDS)


class CodeExecutionResponse(BaseModel):
    success: bool
    output: str
    error: Optional[str] = None
    executionTime: float  # milliseconds
    timestamp: str


# SECURITY: hard caps for the code-execution feature.
MAX_CODE_TIMEOUT_SECONDS = 60        # never let a request monopolize the sandbox
MAX_CODE_LENGTH = 100_000            # bound the submitted source size


@router.post("/code/execute", response_model=CodeExecutionResponse)
async def execute_code(request: CodeExecutionRequest, req: Request):
    """Execute code in a disposable sandbox (rate-limited, size/time-capped)."""
    # Rate limit (same store as file uploads, per-IP).
    client_ip = req.client.host if req.client else "unknown"
    if not check_rate_limit(client_ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded. Please slow down.")

    # Strict language allowlist — reject unknown languages instead of falling back.
    language = (request.language or "").strip().lower()
    if language not in EXEC_LANGUAGES:
        raise HTTPException(status_code=400, detail=f"Unsupported language: {request.language}")

    # Bound input size.
    if len(request.code or "") > MAX_CODE_LENGTH:
        raise HTTPException(status_code=400, detail="Code exceeds maximum allowed length.")

    # Cap the timeout (client may only shorten, never extend).
    timeout = max(1, min(request.timeout, MAX_CODE_TIMEOUT_SECONDS))

    app_logger.info(f"Executing {language} code (timeout: {timeout}s)")

    try:
        # Create sandbox
        sandbox_result = DisposableSandbox.create_sandbox(
            sandbox_name=f"code_exec_{uuid.uuid4().hex[:8]}",
            target_guest_os="auto"
        )

        if not sandbox_result.get("success"):
            raise HTTPException(
                status_code=500,
                detail=f"Failed to create sandbox: {sandbox_result.get('error')}"
            )

        sandbox_id = sandbox_result["sandbox_id"]
        sandbox_dir = Path(sandbox_result["sandbox_path"])

        # Write code to file
        code_file = sandbox_dir / f"main.{_get_extension(language)}"
        code_file.write_text(request.code)

        # Execute code
        start_time = datetime.utcnow()
        exec_result = await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: DisposableSandbox.run_in_sandbox(
                sandbox_id=sandbox_id,
                command=_get_command(language, code_file.name),
                timeout_seconds=timeout
            )
        )
        end_time = datetime.utcnow()

        execution_time = (end_time - start_time).total_seconds() * 1000

        # Cleanup sandbox (always, even on failure)
        try:
            DisposableSandbox.destroy_sandbox(sandbox_id)
        except Exception as e:
            app_logger.warning(f"Sandbox cleanup failed: {e}")

        # Parse result
        success = exec_result.get("success", False)
        output = exec_result.get("stdout", "")
        error = exec_result.get("stderr", "") if not success else None

        return CodeExecutionResponse(
            success=success,
            output=output,
            error=error,
            executionTime=execution_time,
            timestamp=end_time.isoformat()
        )

    except HTTPException:
        raise
    except Exception as e:
        app_logger.error(f"Code execution failed: {e}")
        raise HTTPException(status_code=500, detail=f"Code execution failed: {str(e)}")


# SECURITY: explicit language allowlist for code execution.
EXEC_LANGUAGES = {"python", "javascript", "typescript", "bash", "json", "yaml", "markdown", "text"}


def _get_extension(language: str) -> str:
    """Get file extension for language."""
    extensions = {
        "python": "py",
        "javascript": "js",
        "typescript": "ts",
        "bash": "sh",
        "json": "json",
        "yaml": "yaml",
        "markdown": "md",
        "text": "txt"
    }
    return extensions.get(language.lower(), "txt")


def _get_command(language: str, filename: str) -> str:
    """Get execution command for language."""
    commands = {
        "python": f"python {filename}",
        "javascript": f"node {filename}",
        "typescript": f"ts-node {filename}",
        "bash": f"bash {filename}",
        "json": f"cat {filename}",
        "yaml": f"cat {filename}",
        "markdown": f"cat {filename}",
        "text": f"cat {filename}"
    }
    return commands.get(language.lower(), f"cat {filename}")


# ============================================================================
# Multi-modal Analysis (Phase 6c)
# ============================================================================

class AnalysisRequest(BaseModel):
    fileId: str
    analysisType: str  # "ocr", "vision", "document", "auto"
    promptFocus: Optional[str] = None


class AnalysisResponse(BaseModel):
    success: bool
    type: str
    content: str
    confidence: Optional[float] = None
    metadata: Optional[Dict[str, Any]] = None
    analyzedAt: str


@router.post("/attachments/analyze", response_model=AnalysisResponse)
async def analyze_attachment(request: AnalysisRequest):
    """
    Analyze an attachment (image, document, video, audio, etc.).
    
    Supports automatic analysis type detection based on file category.
    """
    # Find file in metadata
    metadata = load_metadata()
    
    if request.fileId not in metadata:
        raise HTTPException(status_code=404, detail="File not found")
    
    file_info = metadata[request.fileId]
    file_path = Path(file_info['path'])
    
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found on disk")
    
    # Auto-detect analysis type if not specified
    analysis_type = request.analysisType
    if analysis_type == "auto":
        category = file_info.get('category', 'binary')
        if category == 'image':
            analysis_type = 'vision'
        elif category == 'document':
            analysis_type = 'document'
        elif category == 'text' or category == 'code':
            analysis_type = 'document'
        else:
            # For video, audio, archives, etc., just extract metadata
            analysis_type = 'metadata'
    
    app_logger.info(f"Analyzing {request.fileId} with {analysis_type}")
    
    try:
        if analysis_type == "ocr":
            # OCR analysis
            result = await asyncio.get_event_loop().run_in_executor(
                None,
                lambda: OCRReaderTool.extract_text_from_image(str(file_path))
            )
            
            if not result.get("success"):
                raise HTTPException(status_code=500, detail=result.get("error", "OCR failed"))
            
            return AnalysisResponse(
                success=True,
                type="ocr",
                content=result.get("extracted_text", ""),
                confidence=None,
                metadata={"word_count": result.get("word_count", 0)},
                analyzedAt=datetime.utcnow().isoformat()
            )
        
        elif analysis_type == "vision":
            # Vision analysis with LLM
            result = await asyncio.get_event_loop().run_in_executor(
                None,
                lambda: VisionAnalyzerTool.analyze_screen_image(
                    str(file_path),
                    prompt_focus=request.promptFocus
                )
            )
            
            if not result.get("success"):
                raise HTTPException(status_code=500, detail=result.get("error", "Vision analysis failed"))
            
            return AnalysisResponse(
                success=True,
                type="vision",
                content=result.get("ai_analysis", ""),
                confidence=None,
                metadata={
                    "screen_changed": result.get("screen_changed", True),
                    "model_used": result.get("model_used", "unknown")
                },
                analyzedAt=datetime.utcnow().isoformat()
            )
        
        elif analysis_type == "document":
            # Document parsing (PDF, DOCX, text, code)
            ext = file_path.suffix.lower()
            
            if ext == ".pdf":
                content = await _parse_pdf(file_path)
            elif ext in [".doc", ".docx"]:
                content = await _parse_docx(file_path)
            elif ext in [".xls", ".xlsx"]:
                content = await _parse_spreadsheet(file_path)
            elif ext in [".ppt", ".pptx"]:
                content = await _parse_presentation(file_path)
            else:
                # Plain text, code, markdown, etc.
                content = file_path.read_text(errors='ignore')
            
            return AnalysisResponse(
                success=True,
                type="document",
                content=content,
                confidence=None,
                metadata={"file_type": ext, "size": file_info['size']},
                analyzedAt=datetime.utcnow().isoformat()
            )
        
        elif analysis_type == "metadata":
            # Extract metadata for video, audio, archives, etc.
            metadata_info = await _extract_file_metadata(file_path, file_info)
            
            return AnalysisResponse(
                success=True,
                type="metadata",
                content=f"File: {file_info['name']}\nType: {file_info['type']}\nSize: {file_info['size']} bytes\nCategory: {file_info['category']}",
                confidence=None,
                metadata=metadata_info,
                analyzedAt=datetime.utcnow().isoformat()
            )
        
        else:
            raise HTTPException(
                status_code=400,
                detail=f"Unknown analysis type: {analysis_type}"
            )
    
    except HTTPException:
        raise
    except Exception as e:
        app_logger.error(f"Analysis failed: {e}")
        raise HTTPException(status_code=500, detail=f"Analysis failed: {str(e)}")


async def _parse_pdf(file_path: Path) -> str:
    """Parse PDF file and extract text."""
    try:
        from pypdf import PdfReader
        
        def _extract():
            reader = PdfReader(str(file_path))
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        return await asyncio.get_event_loop().run_in_executor(None, _extract)
    except Exception as e:
        app_logger.error(f"PDF parsing failed: {e}")
        raise HTTPException(status_code=500, detail=f"PDF parsing failed: {str(e)}")


async def _parse_docx(file_path: Path) -> str:
    """Parse DOCX file and extract text."""
    try:
        from docx import Document
        
        def _extract():
            doc = Document(str(file_path))
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text
        
        return await asyncio.get_event_loop().run_in_executor(None, _extract)
    except Exception as e:
        app_logger.error(f"DOCX parsing failed: {e}")
        raise HTTPException(status_code=500, detail=f"DOCX parsing failed: {str(e)}")


async def _parse_spreadsheet(file_path: Path) -> str:
    """Parse Excel spreadsheet and extract text."""
    try:
        import pandas as pd
        
        def _extract():
            df = pd.read_excel(str(file_path))
            return df.to_string()
        
        return await asyncio.get_event_loop().run_in_executor(None, _extract)
    except Exception as e:
        app_logger.error(f"Spreadsheet parsing failed: {e}")
        raise HTTPException(status_code=500, detail=f"Spreadsheet parsing failed: {str(e)}")


async def _parse_presentation(file_path: Path) -> str:
    """Parse PowerPoint presentation and extract text."""
    try:
        from pptx import Presentation
        
        def _extract():
            prs = Presentation(str(file_path))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n---\n\n"
            return text
        
        return await asyncio.get_event_loop().run_in_executor(None, _extract)
    except Exception as e:
        app_logger.error(f"Presentation parsing failed: {e}")
        raise HTTPException(status_code=500, detail=f"Presentation parsing failed: {str(e)}")


async def _extract_file_metadata(file_path: Path, file_info: Dict[str, Any]) -> Dict[str, Any]:
    """Extract metadata for video, audio, archives, etc."""
    metadata = {
        'filename': file_info['name'],
        'size': file_info['size'],
        'type': file_info['type'],
        'category': file_info['category'],
        'extension': file_info['extension'],
        'hash': file_info['hash'],
        'uploaded_at': file_info['uploadedAt'],
    }
    
    # Try to extract additional metadata using ffprobe for media files
    if file_info['category'] in ['video', 'audio']:
        try:
            import subprocess
            result = subprocess.run(
                ['ffprobe', '-v', 'quiet', '-print_format', 'json', '-show_format', '-show_streams', str(file_path)],
                capture_output=True,
                text=True,
                timeout=10
            )
            if result.returncode == 0:
                ffprobe_data = json.loads(result.stdout)
                metadata['media_info'] = ffprobe_data
        except Exception as e:
            app_logger.warning(f"ffprobe failed: {e}")
    
    return metadata
