"""PresentationGenerator tests — deterministic .pptx generation from an outline."""

from app.tools.presentation_generator import PresentationGenerator


def _slides():
    return [
        {"title": "Intro", "bullets": ["hello", "world"]},
        {"title": "Details", "bullets": ["one bullet"]},
    ]


def test_generate_presentation(tmp_path):
    out = tmp_path / "deck.pptx"
    res = PresentationGenerator.generate_presentation("My Deck", _slides(), output_path=str(out))
    assert res["success"] is True
    assert res["slide_count"] == 3  # title + 2 content slides
    assert out.exists()
    assert res["file_size_bytes"] > 0


def test_generated_file_is_valid_pptx(tmp_path):
    out = tmp_path / "deck.pptx"
    res = PresentationGenerator.generate_presentation("Deck", _slides(), output_path=str(out))
    assert res["success"] is True

    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) == 3
    assert prs.slides[0].shapes.title.text == "Deck"
    assert prs.slides[1].shapes.title.text == "Intro"


def test_requires_title(tmp_path):
    assert PresentationGenerator.generate_presentation("", _slides())["success"] is False


def test_requires_slides(tmp_path):
    assert PresentationGenerator.generate_presentation("Deck", [])["success"] is False
    assert PresentationGenerator.generate_presentation("Deck", "notalist")["success"] is False


def test_rejects_slide_without_title(tmp_path):
    assert PresentationGenerator.generate_presentation("Deck", [{"bullets": ["x"]}])["success"] is False


def test_forces_pptx_extension(tmp_path):
    out = tmp_path / "deck.txt"
    res = PresentationGenerator.generate_presentation("Deck", _slides(), output_path=str(out))
    assert res["success"] is True
    assert str(res["output_path"]).endswith(".pptx")


def test_subtitle_and_author(tmp_path):
    out = tmp_path / "deck.pptx"
    res = PresentationGenerator.generate_presentation(
        "Deck", _slides(), output_path=str(out), subtitle="A subtitle", author="Arena",
    )
    assert res["success"] is True
    from pptx import Presentation
    prs = Presentation(str(out))
    assert prs.core_properties.author == "Arena"
    assert prs.core_properties.title == "Deck"
