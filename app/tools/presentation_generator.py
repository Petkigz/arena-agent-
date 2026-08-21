"""Presentation generator — deterministic .pptx slides from an outline.

Pure python-pptx + stdlib, no LLM. Turns a structured outline (list of
{title, bullets}) into a real PowerPoint file. The text comes straight from the
caller; nothing is re-worded, so what you wrote is exactly what renders.

Optional dependency: python-pptx (lazy import; missing install degrades to a
clear error).

Safety model (manifest authoritative): Level 1 (draft — creates a new file).
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional

from app.config import settings
from app.utils.logger import app_logger, audit_logger


class PresentationGenerator:
    OUTPUT_DIR = settings.DATA_DIR / "workspace" / "presentations"

    @classmethod
    def ensure_dir(cls) -> None:
        cls.OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    @classmethod
    def generate_presentation(
        cls,
        title: str,
        slides: List[Dict[str, Any]],
        output_path: Optional[str] = None,
        subtitle: str = "",
        author: str = "Arena",
    ) -> Dict[str, Any]:
        """Generate a .pptx with a title slide + one bullet slide per outline item."""
        title = (title or "").strip()
        if not title:
            return {"success": False, "error": "A title is required."}
        slides = cls._normalize_slides(slides)
        if slides is None:
            return {"success": False, "error": "slides must be a non-empty list of {title, bullets}."}

        try:
            from pptx import Presentation
        except ImportError:
            return {"success": False, "error": "python-pptx is not installed. Run: pip install python-pptx"}

        out = Path(output_path) if output_path else cls.OUTPUT_DIR / f"{cls._slug(title)}.pptx"
        if not out.is_absolute():
            out = settings.BASE_DIR / out
        if out.suffix.lower() != ".pptx":
            out = out.with_suffix(".pptx")

        try:
            out.parent.mkdir(parents=True, exist_ok=True)
            prs = Presentation()

            # Title slide.
            s = prs.slides.add_slide(prs.slide_layouts[0])
            s.shapes.title.text = title
            if subtitle:
                s.placeholders[1].text = subtitle

            # Content slides.
            for item in slides:
                sl = prs.slides.add_slide(prs.slide_layouts[1])
                sl.shapes.title.text = item["title"]
                body = sl.placeholders[1].text_frame
                body.clear()
                for i, bullet in enumerate(item["bullets"]):
                    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    p.text = bullet
                    p.level = 0

            if author:
                prs.core_properties.author = author
            prs.core_properties.title = title

            prs.save(str(out))
            audit_logger.info(f"Generated presentation '{title}' → {out}")
            return {
                "success": True,
                "output_path": str(out),
                "slide_count": 1 + len(slides),
                "file_size_bytes": out.stat().st_size,
            }
        except Exception as e:
            app_logger.warning(f"Presentation generation failed: {e}")
            return {"success": False, "error": f"Presentation generation failed: {e}"}

    # ── helpers ─────────────────────────────────────────────────────────────
    @classmethod
    def _normalize_slides(cls, slides) -> Optional[List[Dict[str, Any]]]:
        if not isinstance(slides, list) or not slides:
            return None
        out: List[Dict[str, Any]] = []
        for raw in slides:
            if not isinstance(raw, dict):
                return None
            t = str(raw.get("title", "")).strip()
            if not t:
                return None
            bullets = raw.get("bullets", [])
            if not isinstance(bullets, list):
                return None
            out.append({"title": t, "bullets": [str(b) for b in bullets]})
        return out

    @staticmethod
    def _slug(title: str) -> str:
        return "".join(c if c.isalnum() else "_" for c in title.lower()).strip("_")[:60] or "presentation"
